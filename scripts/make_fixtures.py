#!/usr/bin/env python3
"""Generate the synthetic test corpus in ``tests/fixtures``.

Everything here is written from scratch so the repository never ships
copyrighted material. Every artefact is byte-for-byte reproducible: PDF
generation runs with ReportLab's ``invariant`` flag, OOXML packages get their
zip member timestamps rewritten to a fixed epoch, and the sample git repository
is committed with pinned author/committer dates.

Run it with the ``fixtures`` extra installed::

    uv sync --extra dev --extra fixtures
    uv run python scripts/make_fixtures.py

``--check`` regenerates into a temporary directory and compares hashes against
the committed corpus, which is how CI proves determinism.

Alongside the fixtures the script writes ``ground_truth.json``: the structural
facts each fixture is supposed to carry (expected headings, table cell counts,
article word counts). Conversion tests assert against those facts rather than
against brittle byte equality.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import zipfile
from datetime import UTC, datetime
from pathlib import Path
from typing import Any, Final

REPO_ROOT: Final = Path(__file__).resolve().parent.parent
DEFAULT_OUT: Final = REPO_ROOT / "tests" / "fixtures"

#: Fixed timestamp baked into every generated artefact (2026-01-01T00:00:00Z).
FIXED_DATETIME: Final = datetime(2026, 1, 1, 0, 0, 0, tzinfo=UTC)
FIXED_ZIP_DATE: Final = (2026, 1, 1, 0, 0, 0)
FIXED_GIT_DATE: Final = "2026-01-01T00:00:00+00:00"

#: The article body reused by ``article.html`` and ``boilerplate.html`` so the
#: two fixtures differ *only* in the surrounding boilerplate. That is what makes
#: the Phase 3 boilerplate-reduction measurement meaningful.
ARTICLE_TITLE: Final = "Why Your Context Window Is Mostly Navigation Menus"

ARTICLE_PARAGRAPHS: Final = [
    (
        "Every time a language model reads a web page, it pays for the whole page. "
        "The header pays. The cookie banner pays. The three-deep navigation menu "
        "pays, and so does the newsletter modal, the social share rail, and the "
        "footer that repeats the navigation menu a second time in smaller type."
    ),
    (
        "The article itself is often a small minority of the bytes. On a typical "
        "content site the prose a reader actually came for accounts for a fraction "
        "of the delivered markup, and the rest is scaffolding that carries no "
        "meaning for a model trying to answer a question about the text."
    ),
    (
        "This is why converting a page to Markdown appears to save so many tokens. "
        "The saving is real, but the mechanism is widely misattributed. Markdown "
        "syntax is not meaningfully cheaper than the equivalent HTML tags on a "
        "per-element basis. What actually happens is that the extraction step "
        "throws away the scaffolding, and the format conversion gets the credit."
    ),
    (
        "The distinction matters because it tells you where to spend effort. If the "
        "win came from syntax, you would optimise the serialiser. Because the win "
        "comes from extraction, you should optimise the boilerplate detector, and "
        "you should measure extraction quality rather than output length."
    ),
    (
        "It also tells you when the win will not appear. A page that is already "
        "mostly prose, such as a plain-text RFC or a documentation page rendered "
        "without a shell, has little scaffolding to remove. Running it through an "
        "extractor produces a small reduction and occasionally a regression, "
        "because aggressive extractors sometimes discard content they mistake for "
        "furniture."
    ),
    (
        "There is a second, quieter effect that pushes in the opposite direction. "
        "Structure carries meaning. Headings tell a model how a document is "
        "organised; list markers tell it that items are peers; table pipes tell it "
        "which cell belongs to which column. Strip all of that in pursuit of a "
        "lower token count and you can end up with a cheaper prompt that produces "
        "worse answers."
    ),
    (
        "The practical rule that falls out of this is short enough to remember. "
        "Strip the boilerplate, because it is pure cost. Keep the structure, "
        "because it is load-bearing. And measure both the token count and whether "
        "the answer is still correct, because a token count on its own is not a "
        "quality metric and was never meant to be one."
    ),
]

ARTICLE_SECTIONS: Final = [
    ("Where the tokens actually go", ARTICLE_PARAGRAPHS[0:2]),
    ("The misattributed win", ARTICLE_PARAGRAPHS[2:4]),
    ("When extraction does not help", ARTICLE_PARAGRAPHS[4:6]),
    ("A rule worth remembering", ARTICLE_PARAGRAPHS[6:7]),
]

#: Strings that appear only in the boilerplate. Extraction tests assert these
#: are gone from the converted Markdown.
#: ``long_context.md`` sizing. 42 passages x 6 restatements lands near 80k
#: characters, i.e. roughly 20k tokens for English prose under a GPT-family BPE.
#: The plan asks for "~20k tokens of redundant prose for compression tests".
LONG_CONTEXT_PASSAGES: Final = 42
LONG_CONTEXT_RESTATEMENTS: Final = 6

BOILERPLATE_MARKERS: Final = [
    "Subscribe to our newsletter",
    "SPONSORED: Cut your cloud bill by 40%",
    "Accept all cookies",
    "Trending right now",
    "© 2026 Example Media Group",
    "Follow us on social",
]


# ---------------------------------------------------------------------------
# Determinism helpers
# ---------------------------------------------------------------------------


#: ``dcterms:created`` / ``dcterms:modified`` values inside ``docProps/core.xml``.
_CORE_DATE_RE: Final = re.compile(
    rb"(<dcterms:(?:created|modified)\b[^>]*>)[^<]*(</dcterms:(?:created|modified)>)"
)


def normalise_ooxml(path: Path) -> None:
    """Rewrite an OOXML package so its bytes are reproducible.

    Two sources of non-determinism are removed:

    1. ``python-docx``/``python-pptx``/``openpyxl`` stamp each zip member with
       the current wall-clock time. The archive is rebuilt with a fixed member
       timestamp, a fixed compression level and members sorted by name.
    2. openpyxl overwrites ``dcterms:modified`` with the save time regardless of
       what was set on ``workbook.properties``, so the date fields inside
       ``docProps/core.xml`` are rewritten to the fixed timestamp too.

    Args:
        path: Path to the ``.docx``/``.pptx``/``.xlsx`` file to rewrite in place.
    """
    stamp = FIXED_DATETIME.strftime("%Y-%m-%dT%H:%M:%SZ").encode("ascii")

    with zipfile.ZipFile(path) as original:
        members = sorted(original.infolist(), key=lambda i: i.filename)
        payload = [(info.filename, original.read(info.filename)) for info in members]

    payload = [
        (name, _CORE_DATE_RE.sub(rb"\g<1>" + stamp + rb"\g<2>", data))
        if name == "docProps/core.xml"
        else (name, data)
        for name, data in payload
    ]

    tmp = path.with_suffix(path.suffix + ".tmp")
    with zipfile.ZipFile(tmp, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=9) as rebuilt:
        for name, data in payload:
            info = zipfile.ZipInfo(filename=name, date_time=FIXED_ZIP_DATE)
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = 0o600 << 16
            rebuilt.writestr(info, data)
    tmp.replace(path)


def sha256(path: Path) -> str:
    """Return the hex SHA-256 digest of a file.

    Args:
        path: File to digest.

    Returns:
        Lowercase hex digest.
    """
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for block in iter(lambda: handle.read(1 << 20), b""):
            digest.update(block)
    return digest.hexdigest()


# ---------------------------------------------------------------------------
# PDF fixtures
# ---------------------------------------------------------------------------

TABLE_HEADER: Final = ["Backend", "License", "Runtime", "Tables", "Pages/sec"]
TABLE_ROWS: Final = [
    ["markitdown", "MIT", "CPU", "weak", "12.0"],
    ["docling", "MIT", "CPU", "strong", "0.8"],
    ["pdfplumber", "MIT", "CPU", "good", "3.4"],
    ["pypdf", "BSD-3", "CPU", "none", "18.5"],
    ["pymupdf4llm", "AGPL-3.0", "CPU", "good", "11.1"],
    ["marker", "GPL-3.0", "GPU", "strong", "1.9"],
]


def _configure_reportlab() -> None:
    """Put ReportLab into invariant mode so generated PDFs are reproducible."""
    from reportlab import rl_config

    rl_config.invariant = 1
    rl_config.useA85 = 0


def build_simple_pdf(out: Path) -> dict[str, Any]:
    """Write ``simple.pdf``: a plain multi-page digital PDF with headings.

    Args:
        out: Fixture output directory.

    Returns:
        Ground-truth facts for the fixture.
    """
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer

    _configure_reportlab()
    styles = getSampleStyleSheet()
    path = out / "simple.pdf"

    headings = ["Introduction", "Method", "Results", "Discussion"]
    story: list[Any] = [Paragraph(ARTICLE_TITLE, styles["Title"]), Spacer(1, 12)]
    for index, heading in enumerate(headings):
        story.append(Paragraph(heading, styles["Heading1"]))
        for paragraph in ARTICLE_PARAGRAPHS[index * 2 : index * 2 + 2] or ARTICLE_PARAGRAPHS[:1]:
            story.append(Paragraph(paragraph, styles["BodyText"]))
            story.append(Spacer(1, 6))
        if index == 1:
            story.append(PageBreak())

    doc = SimpleDocTemplate(
        str(path),
        pagesize=LETTER,
        title=ARTICLE_TITLE,
        author="tokenmill fixtures",
        subject="synthetic test corpus",
    )
    doc.build(story)
    return {
        "description": "Baseline digital PDF: title, four H1 sections, multi-page.",
        "expected_headings": [ARTICLE_TITLE, *headings],
        "min_pages": 2,
        "must_contain": ["Strip the boilerplate", "load-bearing"],
    }


def build_tables_pdf(out: Path) -> dict[str, Any]:
    """Write ``tables.pdf``: a single grid table used for table-fidelity checks.

    Args:
        out: Fixture output directory.

    Returns:
        Ground-truth facts for the fixture.
    """
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    _configure_reportlab()
    styles = getSampleStyleSheet()
    path = out / "tables.pdf"

    data = [TABLE_HEADER, *TABLE_ROWS]
    table = Table(data, hAlign="LEFT")
    table.setStyle(
        TableStyle(
            [
                ("GRID", (0, 0), (-1, -1), 0.5, colors.black),
                ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("ALIGN", (4, 1), (4, -1), "RIGHT"),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ]
        )
    )

    story: list[Any] = [
        Paragraph("Converter Comparison", styles["Title"]),
        Spacer(1, 12),
        Paragraph("Backend characteristics", styles["Heading1"]),
        Paragraph(
            "The table below is the fixture's reason for existing: a converter "
            "that flattens it into prose has lost the data.",
            styles["BodyText"],
        ),
        Spacer(1, 12),
        table,
        Spacer(1, 12),
        Paragraph(
            "Figures are illustrative placeholders for structural testing and are "
            "not measurements of any real backend.",
            styles["BodyText"],
        ),
    ]
    SimpleDocTemplate(
        str(path),
        pagesize=LETTER,
        title="Converter Comparison",
        author="tokenmill fixtures",
    ).build(story)

    return {
        "description": "One 7x5 grid table (header + 6 rows) for table-fidelity checks.",
        "expected_headings": ["Converter Comparison", "Backend characteristics"],
        "table_count": 1,
        "table_columns": len(TABLE_HEADER),
        "table_rows_including_header": len(TABLE_ROWS) + 1,
        "table_cells": (len(TABLE_ROWS) + 1) * len(TABLE_HEADER),
        "table_header": list(TABLE_HEADER),
        "table_first_column": [row[0] for row in TABLE_ROWS],
        "must_contain": ["pdfplumber", "AGPL-3.0", "18.5"],
    }


def build_twocolumn_pdf(out: Path) -> dict[str, Any]:
    """Write ``twocolumn.pdf``: two-column body text for reading-order checks.

    Args:
        out: Fixture output directory.

    Returns:
        Ground-truth facts for the fixture.
    """
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import BaseDocTemplate, Frame, PageTemplate, Paragraph, Spacer

    _configure_reportlab()
    styles = getSampleStyleSheet()
    path = out / "twocolumn.pdf"

    page_width, page_height = LETTER
    margin = 54.0
    gutter = 18.0
    column_width = (page_width - 2 * margin - gutter) / 2
    frame_height = page_height - 2 * margin

    left = Frame(margin, margin, column_width, frame_height, id="left")
    right = Frame(margin + column_width + gutter, margin, column_width, frame_height, id="right")

    doc = BaseDocTemplate(
        str(path),
        pagesize=LETTER,
        title="Two Column Reading Order",
        author="tokenmill fixtures",
    )
    doc.addPageTemplates([PageTemplate(id="twocol", frames=[left, right])])

    # Sentinel sentences let a test assert the reading order came out as
    # left-column-then-right-column rather than interleaved by scan line.
    ordered_markers = [f"ORDERMARK {n:02d}." for n in range(1, 13)]
    story: list[Any] = [Paragraph("Two Column Reading Order", styles["Heading1"])]
    for index, marker in enumerate(ordered_markers):
        body = ARTICLE_PARAGRAPHS[index % len(ARTICLE_PARAGRAPHS)]
        story.append(Paragraph(f"{marker} {body}", styles["BodyText"]))
        story.append(Spacer(1, 6))
    doc.build(story)

    return {
        "description": ("Two-column layout; ORDERMARK sentinels encode correct reading order."),
        "expected_headings": ["Two Column Reading Order"],
        "order_markers": [marker.rstrip(".") for marker in ordered_markers],
        "must_contain": ["ORDERMARK 01", "ORDERMARK 12"],
    }


def build_scanned_pdf(out: Path) -> dict[str, Any]:
    """Write ``scanned.pdf`` by rasterising ``simple.pdf`` into page images.

    The result has no text layer at all, which is exactly what the OCR path
    needs to be exercised against, and what non-OCR backends should report as an
    empty or near-empty extraction.

    Args:
        out: Fixture output directory.

    Returns:
        Ground-truth facts for the fixture.

    Raises:
        FileNotFoundError: If ``simple.pdf`` has not been generated yet.
    """
    import pypdfium2 as pdfium
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas as pdfcanvas

    _configure_reportlab()
    source = out / "simple.pdf"
    if not source.exists():
        msg = f"scanned.pdf needs {source} to exist first"
        raise FileNotFoundError(msg)

    path = out / "scanned.pdf"
    scale = 150 / 72  # render at 150 DPI
    document = pdfium.PdfDocument(str(source))
    try:
        images = [
            document[index].render(scale=scale, grayscale=True).to_pil()
            for index in range(len(document))
        ]
    finally:
        document.close()

    page_width, page_height = LETTER
    canvas = pdfcanvas.Canvas(str(path), pagesize=LETTER)
    canvas.setTitle("Scanned Document")
    canvas.setAuthor("tokenmill fixtures")
    for image in images:
        canvas.drawImage(
            ImageReader(image), 0, 0, width=page_width, height=page_height, preserveAspectRatio=True
        )
        canvas.showPage()
    canvas.save()

    return {
        "description": "Rasterised copy of simple.pdf: images only, no text layer.",
        "page_count": len(images),
        "has_text_layer": False,
        "requires_ocr": True,
        "source_fixture": "simple.pdf",
    }


def build_corrupt_pdf(out: Path) -> dict[str, Any]:
    """Write ``corrupt.pdf``: a truncated PDF for exercising the error path.

    Args:
        out: Fixture output directory.

    Returns:
        Ground-truth facts for the fixture.

    Raises:
        FileNotFoundError: If ``simple.pdf`` has not been generated yet.
    """
    source = out / "simple.pdf"
    if not source.exists():
        msg = f"corrupt.pdf needs {source} to exist first"
        raise FileNotFoundError(msg)

    path = out / "corrupt.pdf"
    data = source.read_bytes()
    # Keep the %PDF- header so sniffers still identify it as a PDF, then cut it
    # off well before the xref table.
    path.write_bytes(data[: len(data) // 3])

    return {
        "description": (
            "Truncated PDF (first third of simple.pdf). Every backend must fail cleanly."
        ),
        "expect_conversion_error": True,
        "keeps_pdf_magic": True,
    }


# ---------------------------------------------------------------------------
# Office fixtures
# ---------------------------------------------------------------------------


def _stamp_docx_properties(document: Any) -> None:
    """Pin a python-docx document's core properties to the fixed timestamp.

    Args:
        document: The ``docx.Document`` instance to stamp.
    """
    properties = document.core_properties
    properties.author = "tokenmill fixtures"
    properties.created = FIXED_DATETIME.replace(tzinfo=None)
    properties.modified = FIXED_DATETIME.replace(tzinfo=None)
    properties.last_modified_by = "tokenmill fixtures"
    properties.revision = 1


def build_report_docx(out: Path) -> dict[str, Any]:
    """Write ``report.docx``: nested headings, nested lists, a table and a footnote.

    Args:
        out: Fixture output directory.

    Returns:
        Ground-truth facts for the fixture.
    """
    from docx import Document

    path = out / "report.docx"
    document = Document()

    document.add_heading("Context Efficiency Report", level=0)
    document.add_paragraph(ARTICLE_PARAGRAPHS[0])

    headings: list[tuple[str, int]] = []
    for section_title, paragraphs in ARTICLE_SECTIONS:
        document.add_heading(section_title, level=1)
        headings.append((section_title, 1))
        for paragraph in paragraphs:
            document.add_paragraph(paragraph)
        subsection = f"{section_title}: detail"
        document.add_heading(subsection, level=2)
        headings.append((subsection, 2))
        document.add_paragraph("Supporting notes for the section above.")

    document.add_heading("Checklist", level=1)
    headings.append(("Checklist", 1))
    for item in ["Strip navigation", "Strip advertising", "Strip cookie banners"]:
        document.add_paragraph(item, style="List Bullet")
    for item in ["Keep headings", "Keep list markers", "Keep table structure"]:
        document.add_paragraph(item, style="List Number")
    document.add_paragraph("Nested detail under the last item", style="List Number 2")

    document.add_heading("Measurements", level=1)
    headings.append(("Measurements", 1))
    table = document.add_table(rows=1, cols=3)
    table.style = "Table Grid"
    for cell, text in zip(table.rows[0].cells, ["Stage", "Tokens", "Delta"], strict=True):
        cell.text = text
    measurement_rows = [
        ["source", "16180", "-"],
        ["converted", "3150", "-80.5%"],
        ["post-processed", "2980", "-81.6%"],
    ]
    for row_values in measurement_rows:
        cells = table.add_row().cells
        for cell, text in zip(cells, row_values, strict=True):
            cell.text = text
    document.add_paragraph(
        "Note: figures above are illustrative placeholders for structural testing, "
        "not measurements produced by this project."
    )

    _stamp_docx_properties(document)
    document.save(str(path))
    normalise_ooxml(path)

    return {
        "description": (
            "Deep heading hierarchy (H0-H2), bullet + numbered + nested lists, 4x3 table."
        ),
        "expected_headings": [("Context Efficiency Report", 0), *headings],
        "heading_levels_present": [0, 1, 2],
        "table_count": 1,
        "table_columns": 3,
        "table_rows_including_header": len(measurement_rows) + 1,
        "bullet_items": ["Strip navigation", "Strip advertising", "Strip cookie banners"],
        "numbered_items": ["Keep headings", "Keep list markers", "Keep table structure"],
        "must_contain": ["Nested detail under the last item", "16180"],
    }


def build_unicode_docx(out: Path) -> dict[str, Any]:
    """Write ``unicode.docx``: Urdu, Arabic, CJK, Cyrillic, Devanagari and emoji.

    Args:
        out: Fixture output directory.

    Returns:
        Ground-truth facts for the fixture.
    """
    from docx import Document

    path = out / "unicode.docx"
    samples: list[tuple[str, str]] = [
        ("Urdu", "ٹوکن کی بچت ایک اہم مسئلہ ہے۔ دستاویز کو مارک ڈاؤن میں تبدیل کریں۔"),
        ("Arabic", "توفير الرموز مسألة مهمة. حوّل المستند إلى ماركداون."),
        ("Chinese", "令牌节省是一个重要问题。请将文档转换为 Markdown 格式。"),
        ("Japanese", "トークンの節約は重要な課題です。ドキュメントをマークダウンに変換します。"),
        ("Korean", "토큰 절약은 중요한 문제입니다. 문서를 마크다운으로 변환하십시오."),
        ("Russian", "Экономия токенов — важная задача. Преобразуйте документ в Markdown."),
        ("Hindi", "टोकन की बचत एक महत्वपूर्ण समस्या है। दस्तावेज़ को मार्कडाउन में बदलें।"),
        ("Greek", "Η εξοικονόμηση διακριτικών είναι σημαντική. Μετατρέψτε το έγγραφο."),
        ("Emoji", "Tokens saved 🎉 — pipeline 📄 ➡️ 📝 at 80% 🔻 (family: 👨‍👩‍👧‍👦, flag: 🇵🇰)"),
        ("Math", "Compression ratio ρ ≈ 0.31, Δtokens = −13 030, ∑ᵢ tᵢ ≤ 4 096"),
    ]

    document = Document()
    document.add_heading("Unicode Round-Trip", level=0)
    document.add_paragraph(
        "Each row below must survive conversion byte-for-byte in its decoded form."
    )
    for label, text in samples:
        document.add_heading(label, level=1)
        document.add_paragraph(text)

    _stamp_docx_properties(document)
    document.save(str(path))
    normalise_ooxml(path)

    return {
        "description": (
            "Multi-script text (RTL, CJK, Devanagari, emoji ZWJ, math) for encoding regressions."
        ),
        "expected_headings": ["Unicode Round-Trip", *[label for label, _ in samples]],
        "scripts": dict(samples),
        "must_contain": [text for _, text in samples],
    }


def build_deck_pptx(out: Path) -> dict[str, Any]:
    """Write ``deck.pptx``: titled slides with bullets and speaker notes.

    Args:
        out: Fixture output directory.

    Returns:
        Ground-truth facts for the fixture.
    """
    from pptx import Presentation
    from pptx.util import Inches

    path = out / "deck.pptx"
    presentation = Presentation()
    title_layout = presentation.slide_layouts[0]
    bullet_layout = presentation.slide_layouts[1]

    slides_spec: list[tuple[str, list[str], str]] = [
        (
            "Where Your Tokens Go",
            ["Navigation", "Advertising", "Cookie banners", "The article you wanted"],
            "Open by asking the room to guess the split. It is always worse than they think.",
        ),
        (
            "Strip Boilerplate, Keep Structure",
            ["Boilerplate is pure cost", "Headings carry meaning", "Measure both axes"],
            "This is the one slide people should photograph.",
        ),
        (
            "Measuring Honestly",
            ["Tokens before", "Tokens after", "Fidelity against ground truth"],
            "A token count without a fidelity score is not a result.",
        ),
    ]

    title_slide = presentation.slides.add_slide(title_layout)
    title_slide.shapes.title.text = "Context Efficiency"
    title_slide.placeholders[1].text = "A short deck fixture for tokenmill"
    title_slide.notes_slide.notes_text_frame.text = "Title slide speaker note."

    for slide_title, bullets, note in slides_spec:
        slide = presentation.slides.add_slide(bullet_layout)
        slide.shapes.title.text = slide_title
        body = slide.placeholders[1].text_frame
        body.text = bullets[0]
        for bullet in bullets[1:]:
            paragraph = body.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 1
        slide.notes_slide.notes_text_frame.text = note

    blank = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = blank.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    textbox.text_frame.text = "Appendix: raw numbers live in the repository."

    presentation.core_properties.author = "tokenmill fixtures"
    presentation.core_properties.created = FIXED_DATETIME.replace(tzinfo=None)
    presentation.core_properties.modified = FIXED_DATETIME.replace(tzinfo=None)
    presentation.core_properties.last_modified_by = "tokenmill fixtures"
    presentation.core_properties.revision = 1
    presentation.save(str(path))
    normalise_ooxml(path)

    return {
        "description": (
            "5 slides: title, three bulleted slides with speaker notes, one blank+textbox."
        ),
        "slide_count": 5,
        "slide_titles": ["Context Efficiency", *[title for title, _, _ in slides_spec]],
        "speaker_notes": [
            "Title slide speaker note.",
            *[note for _, _, note in slides_spec],
        ],
        "must_contain": ["Cookie banners", "Appendix: raw numbers live in the repository."],
    }


def build_data_xlsx(out: Path) -> dict[str, Any]:
    """Write ``data.xlsx``: three sheets of tabular data including a formula.

    Args:
        out: Fixture output directory.

    Returns:
        Ground-truth facts for the fixture.
    """
    from openpyxl import Workbook

    path = out / "data.xlsx"
    workbook = Workbook()

    backends = workbook.active
    backends.title = "backends"
    backends.append(list(TABLE_HEADER))
    for row in TABLE_ROWS:
        backends.append([row[0], row[1], row[2], row[3], float(row[4])])

    corpus = workbook.create_sheet("corpus")
    corpus.append(["fixture", "format", "bytes", "pages"])
    corpus_rows: list[list[object]] = [
        ["simple.pdf", "pdf", 4200, 2],
        ["tables.pdf", "pdf", 3100, 1],
        ["report.docx", "docx", 38000, 0],
        ["deck.pptx", "pptx", 29000, 0],
    ]
    for corpus_row in corpus_rows:
        corpus.append(corpus_row)

    totals = workbook.create_sheet("totals")
    totals.append(["metric", "value"])
    totals.append(["backend_count", len(TABLE_ROWS)])
    totals.append(["corpus_items", len(corpus_rows)])
    totals.append(["mean_pages_per_sec", "=AVERAGE(backends!E2:E7)"])

    workbook.properties.creator = "tokenmill fixtures"
    workbook.properties.lastModifiedBy = "tokenmill fixtures"
    workbook.properties.created = FIXED_DATETIME.replace(tzinfo=None)
    workbook.properties.modified = FIXED_DATETIME.replace(tzinfo=None)
    workbook.save(str(path))
    normalise_ooxml(path)

    return {
        "description": (
            "Three sheets (backends, corpus, totals); totals holds one AVERAGE formula."
        ),
        "sheet_names": ["backends", "corpus", "totals"],
        "sheet_row_counts": {
            "backends": len(TABLE_ROWS) + 1,
            "corpus": len(corpus_rows) + 1,
            "totals": 4,
        },
        "has_formula": True,
        "must_contain": ["pdfplumber", "simple.pdf", "backend_count"],
    }


# ---------------------------------------------------------------------------
# HTML fixtures
# ---------------------------------------------------------------------------


def _article_body_html() -> str:
    """Render the shared article body as HTML.

    Returns:
        The ``<article>`` element markup used by both HTML fixtures.
    """
    parts = [
        "  <article>",
        f"    <h1>{ARTICLE_TITLE}</h1>",
        '    <p class="standfirst">A short look at where context windows actually go.</p>',
    ]
    for section_title, paragraphs in ARTICLE_SECTIONS:
        parts.append(f"    <h2>{section_title}</h2>")
        parts.extend(f"    <p>{paragraph}</p>" for paragraph in paragraphs)
    parts.append("    <h2>Summary table</h2>")
    parts.append("    <table>")
    parts.append(
        "      <thead><tr>" + "".join(f"<th>{h}</th>" for h in TABLE_HEADER) + "</tr></thead>"
    )
    parts.append("      <tbody>")
    for row in TABLE_ROWS:
        parts.append("        <tr>" + "".join(f"<td>{cell}</td>" for cell in row) + "</tr>")
    parts.append("      </tbody>")
    parts.append("    </table>")
    parts.append("  </article>")
    return "\n".join(parts)


def build_article_html(out: Path) -> dict[str, Any]:
    """Write ``article.html``: the clean extraction baseline.

    Args:
        out: Fixture output directory.

    Returns:
        Ground-truth facts for the fixture.
    """
    path = out / "article.html"
    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<title>{ARTICLE_TITLE}</title>
</head>
<body>
{_article_body_html()}
</body>
</html>
"""
    path.write_text(html, encoding="utf-8", newline="\n")
    return _html_ground_truth(
        "Clean article with no boilerplate: the extraction baseline.",
        has_boilerplate=False,
    )


def build_boilerplate_html(out: Path) -> dict[str, Any]:
    """Write ``boilerplate.html``: the same article buried in nav, ads and scripts.

    Args:
        out: Fixture output directory.

    Returns:
        Ground-truth facts for the fixture.
    """
    path = out / "boilerplate.html"

    nav_sections = {
        "News": ["World", "Politics", "Business", "Technology", "Science", "Health"],
        "Opinion": ["Editorials", "Columnists", "Letters", "Guest Essays"],
        "Culture": ["Books", "Film", "Music", "Television", "Theatre", "Art"],
        "Lifestyle": ["Food", "Travel", "Style", "Wellbeing", "Home", "Gardening"],
        "More": ["Newsletters", "Podcasts", "Video", "Crosswords", "Jobs", "Archive"],
    }

    def slug(value: str) -> str:
        return value.lower().replace(" ", "-")

    nav_items = "\n".join(
        "    <li><strong>{}</strong><ul>{}</ul></li>".format(
            section,
            "".join(
                f'<li><a href="/{slug(section)}/{slug(link)}">{link}</a></li>' for link in links
            ),
        )
        for section, links in nav_sections.items()
    )

    trending = "\n".join(
        f'      <li><a href="/trending/{n}">Ten things about topic number {n} '
        f"that will change how you work forever</a></li>"
        for n in range(1, 13)
    )

    inline_script = (
        "(function(){var q=window.__adq=window.__adq||[];"
        "for(var i=0;i<40;i++){q.push({slot:'slot-'+i,sizes:[[300,250],[728,90]],"
        "targeting:{section:'technology',page:'article',variant:'control'}});}"
        "window.dataLayer=window.dataLayer||[];"
        "window.dataLayer.push({event:'pageview',contentGroup:'technology',"
        "author:'staff',wordCount:1200,paywall:false});})();"
    )
    inline_style = (
        ".nav{display:flex;gap:1rem}.nav a{color:#333;text-decoration:none}"
        ".ad{min-height:250px;background:#eee;border:1px dashed #ccc}"
        ".cookie{position:fixed;bottom:0;left:0;right:0;background:#111;color:#fff}"
        ".modal{position:fixed;inset:0;background:rgba(0,0,0,.6)}"
        ".footer-links{columns:4}.social{display:flex}"
    ) * 6

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{ARTICLE_TITLE} | Example Media Group</title>
<meta name="description" content="A short look at where context windows actually go.">
<style>{inline_style}</style>
<script>{inline_script}</script>
</head>
<body>
<div class="cookie" id="cookie-banner">
  <p>We and our 847 partners use cookies and similar technologies to store and
  access information on your device, personalise advertising, and measure
  audience engagement. You may change your preferences at any time.</p>
  <button>Accept all cookies</button><button>Manage preferences</button><button>Reject all</button>
</div>

<header>
  <a href="/" class="logo">Example Media Group</a>
  <nav class="nav">
  <ul>
{nav_items}
  </ul>
  </nav>
  <form action="/search"><input type="search" name="q" placeholder="Search"><button>Go</button></form>
  <a class="cta" href="/subscribe">Subscribe from $1/week</a>
</header>

<div class="ad" id="ad-leaderboard">
  <p>SPONSORED: Cut your cloud bill by 40% with one weird configuration change.</p>
  <a href="https://ads.example.com/click?id=1">Learn more</a>
</div>

<main>
{_article_body_html()}

  <aside class="ad" id="ad-mpu">
    <p>SPONSORED: Cut your cloud bill by 40% with one weird configuration change.</p>
  </aside>

  <aside class="social">
    <p>Follow us on social</p>
    <ul><li><a href="#">Share</a></li><li><a href="#">Post</a></li><li><a href="#">Send</a></li></ul>
  </aside>

  <aside class="trending">
    <h3>Trending right now</h3>
    <ul>
{trending}
    </ul>
  </aside>

  <div class="modal" id="newsletter-modal">
    <h3>Subscribe to our newsletter</h3>
    <p>Get the stories that matter, every weekday morning, straight to your inbox.</p>
    <form><input type="email" placeholder="you@example.com"><button>Sign up</button></form>
  </div>
</main>

<footer>
  <div class="footer-links">
  <ul>
{nav_items}
  </ul>
  </div>
  <p>© 2026 Example Media Group. All rights reserved.
  Terms of Service · Privacy Policy · Cookie Policy · Accessibility · Contact ·
  Careers · Advertise with us · Corrections · Complaints</p>
</footer>
<script>{inline_script}</script>
</body>
</html>
"""
    path.write_text(html, encoding="utf-8", newline="\n")

    truth = _html_ground_truth(
        "The article.html body wrapped in heavy nav, ads, scripts, modals and footer.",
        has_boilerplate=True,
    )
    truth["boilerplate_markers_must_be_absent"] = list(BOILERPLATE_MARKERS)
    truth["paired_clean_fixture"] = "article.html"
    return truth


def _html_ground_truth(description: str, *, has_boilerplate: bool) -> dict[str, Any]:
    """Build the shared ground-truth block for the two HTML fixtures.

    Args:
        description: Human-readable description of the fixture.
        has_boilerplate: Whether the fixture wraps the article in boilerplate.

    Returns:
        Ground-truth facts for the fixture.
    """
    body_words = sum(len(paragraph.split()) for paragraph in ARTICLE_PARAGRAPHS)
    # Counts the <p> body only: not the standfirst, headings or table cells.
    return {
        "description": description,
        "has_boilerplate": has_boilerplate,
        "article_title": ARTICLE_TITLE,
        "expected_headings": [ARTICLE_TITLE]
        + [title for title, _ in ARTICLE_SECTIONS]
        + ["Summary table"],
        "article_paragraph_count": len(ARTICLE_PARAGRAPHS),
        "article_body_word_count": body_words,
        "table_count": 1,
        "table_columns": len(TABLE_HEADER),
        "table_rows_including_header": len(TABLE_ROWS) + 1,
        "must_contain": ["Strip the boilerplate", "load-bearing", "pdfplumber"],
    }


# ---------------------------------------------------------------------------
# Text and repository fixtures
# ---------------------------------------------------------------------------


def build_long_context_md(out: Path) -> dict[str, Any]:
    """Write ``long_context.md``: deliberately redundant prose for compression tests.

    Args:
        out: Fixture output directory.

    Returns:
        Ground-truth facts for the fixture.
    """
    path = out / "long_context.md"

    # A needle the compression tests can check survived. Prompt compressors are
    # supposed to keep salient facts and discard restatement; this fixture is
    # mostly restatement with exactly one specific, checkable fact.
    needle = "The internal build identifier for this corpus is RSD-TOKENMILL-4417."

    lines: list[str] = [
        "# Retrieval Context Bundle",
        "",
        "This document is intentionally redundant. It exists so that prompt",
        "compression can be measured against text whose information density is",
        "known to be low.",
        "",
        needle,
        "",
    ]

    restatements = [
        "The retrieved passage restates that boilerplate removal accounts for most "
        "of the observed token reduction when documents are converted to Markdown.",
        "As noted previously, and as will be noted again, boilerplate removal is the "
        "dominant term and Markdown syntax itself is close to a rounding error.",
        "It is worth repeating that navigation menus, cookie banners and advertising "
        "slots contribute a large share of the delivered bytes on content sites.",
        "To restate the point in different words: the scaffolding around an article "
        "is expensive, and the article itself is comparatively cheap.",
        "Once more for emphasis, structure such as headings and list markers carries "
        "meaning and should be preserved even when it costs a small number of tokens.",
        "Repeating the caveat: a lower token count is not by itself evidence of a "
        "better prompt, because the answer may have become worse.",
    ]

    for chapter in range(1, LONG_CONTEXT_PASSAGES + 1):
        lines.append(f"## Passage {chapter:02d}")
        lines.append("")
        for repeat in range(LONG_CONTEXT_RESTATEMENTS):
            sentence = restatements[(chapter + repeat) % len(restatements)]
            lines.append(
                f"{sentence} This is passage {chapter:02d}, restatement "
                f"{repeat + 1} of {LONG_CONTEXT_RESTATEMENTS}, "
                f"and it adds no information beyond what the earlier passages already "
                f"established about boilerplate, structure and measurement."
            )
            lines.append("")

    lines.append("## Conclusion")
    lines.append("")
    lines.append(needle)
    lines.append("")
    text = "\n".join(lines)
    path.write_text(text, encoding="utf-8", newline="\n")

    return {
        "description": (
            "Deliberately redundant prose sized for prompt-compression tests, "
            "with one checkable needle fact."
        ),
        "needle": needle,
        "needle_occurrences": 2,
        "word_count": len(text.split()),
        "character_count": len(text),
        "passage_count": LONG_CONTEXT_PASSAGES,
        "restatements_per_passage": LONG_CONTEXT_RESTATEMENTS,
        # Deliberately NOT a token count. Token counts are tokenizer-dependent
        # and this repository does not publish numbers it has not measured; the
        # size target below was hit with the ~4-characters-per-token rule of
        # thumb, and Phase 1 records real measured counts once a tokenizer is
        # wired up. See PROGRESS.md.
        "token_count": None,
        "token_count_note": (
            "Not measured yet. Sized to roughly 20k tokens using the ~4 chars/token "
            "rule of thumb for English prose; Phase 1 replaces this with a measured "
            "count per tokenizer."
        ),
    }


SAMPLE_REPO_FILES: Final[dict[str, str]] = {
    "README.md": """# widgetlib

A tiny synthetic package used as a fixture for repository ingestion.

## Install

```bash
pip install widgetlib
```

## Usage

```python
from widgetlib import Widget

print(Widget("demo").render())
```
""",
    ".gitignore": """__pycache__/
*.pyc
build/
dist/
.venv/
secrets.env
""",
    "pyproject.toml": """[project]
name = "widgetlib"
version = "0.1.0"
requires-python = ">=3.11"
""",
    "src/widgetlib/__init__.py": '''"""A tiny synthetic package."""

from widgetlib.core import Widget

__all__ = ["Widget"]
__version__ = "0.1.0"
''',
    "src/widgetlib/core.py": '''"""Core widget implementation."""

from __future__ import annotations

from dataclasses import dataclass


@dataclass(frozen=True)
class Widget:
    """A widget with a name and an optional size."""

    name: str
    size: int = 1

    def render(self) -> str:
        """Return the widget's textual representation."""
        return f"<{self.name} size={self.size}>"

    def scaled(self, factor: int) -> Widget:
        """Return a copy of this widget scaled by ``factor``."""
        if factor <= 0:
            raise ValueError("factor must be positive")
        return Widget(self.name, self.size * factor)
''',
    "src/widgetlib/utils.py": '''"""Helpers that are deliberately unremarkable."""

from __future__ import annotations

from collections.abc import Iterable


def total_size(widgets: Iterable[object]) -> int:
    """Sum the ``size`` attribute of every widget in ``widgets``."""
    return sum(getattr(widget, "size", 0) for widget in widgets)
''',
    "tests/test_core.py": '''"""Tests for widgetlib.core."""

from widgetlib import Widget


def test_render() -> None:
    assert Widget("demo").render() == "<demo size=1>"


def test_scaled() -> None:
    assert Widget("demo", 2).scaled(3).size == 6
''',
    "docs/design.md": """# Design notes

widgetlib exists only so that repository-ingestion backends have something with
a directory tree, a README, source, tests, docs, an ignored file and a binary
blob to walk.
""",
}

#: Files that exist on disk but must be excluded from ingestion output.
SAMPLE_REPO_IGNORED: Final[dict[str, str]] = {
    "secrets.env": "API_KEY=this-should-never-appear-in-ingested-output\n",
}


def ensure_sample_repo_git(root: Path) -> str:
    """Make ``root`` a git repository with the fixture's pinned commit.

    Idempotent: if the repository already exists, its HEAD is returned
    unchanged. Author and committer dates are pinned, so a repository recreated
    from the committed working files always lands on the same commit hash.

    This is separate from writing the files because the ``.git`` directory is
    deliberately **not** committed. Git treats a nested ``.git`` as a submodule
    boundary and stores a gitlink instead of the contents, which would leave
    anyone cloning tokenmill with an empty ``sample_repo/``. The working files
    are committed; the repository is materialised here on demand.

    Args:
        root: The ``sample_repo`` directory.

    Returns:
        The 40-character hex SHA of HEAD.

    Raises:
        RuntimeError: If ``git`` is not on PATH.
    """
    git = shutil.which("git")
    if git is None:
        msg = "git is required to build the sample_repo fixture"
        raise RuntimeError(msg)

    env = {
        **os.environ,
        "GIT_AUTHOR_NAME": "tokenmill fixtures",
        "GIT_AUTHOR_EMAIL": "fixtures@example.invalid",
        "GIT_AUTHOR_DATE": FIXED_GIT_DATE,
        "GIT_COMMITTER_NAME": "tokenmill fixtures",
        "GIT_COMMITTER_EMAIL": "fixtures@example.invalid",
        "GIT_COMMITTER_DATE": FIXED_GIT_DATE,
        "GIT_CONFIG_GLOBAL": os.devnull,
        "GIT_CONFIG_SYSTEM": os.devnull,
    }

    def run(*args: str) -> str:
        # Fixed argv list, no shell, absolute git path.
        return subprocess.run(  # noqa: S603
            [git, *args],
            cwd=root,
            env=env,
            check=True,
            capture_output=True,
            text=True,
        ).stdout.strip()

    if not (root / ".git").is_dir():
        run("init", "--quiet", "--initial-branch=main")
        run("add", "--all")
        run("commit", "--quiet", "-m", "Initial commit")
    return run("rev-parse", "HEAD")


def build_sample_repo(out: Path) -> dict[str, Any]:
    """Write ``sample_repo/``: a small git repository with a pinned commit.

    Args:
        out: Fixture output directory.

    Returns:
        Ground-truth facts for the fixture.
    """
    root = out / "sample_repo"
    if root.exists():
        shutil.rmtree(root)
    root.mkdir(parents=True)

    for relative, content in {**SAMPLE_REPO_FILES, **SAMPLE_REPO_IGNORED}.items():
        target = root / relative
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_text(content, encoding="utf-8", newline="\n")

    # A small binary blob: ingestion backends must skip it, not mangle it.
    binary = root / "assets" / "logo.bin"
    binary.parent.mkdir(parents=True, exist_ok=True)
    binary.write_bytes(bytes(range(256)) * 8)

    head = ensure_sample_repo_git(root)

    tracked = sorted([*SAMPLE_REPO_FILES, "assets/logo.bin"])
    return {
        "description": (
            "Small git repo: src/, tests/, docs/, a .gitignore'd secret and a binary blob."
        ),
        "tracked_files": tracked,
        "tracked_file_count": len(tracked),
        "ignored_files": sorted(SAMPLE_REPO_IGNORED),
        "binary_files": ["assets/logo.bin"],
        "head_commit": head,
        "default_branch": "main",
        "git_dir_is_committed": False,
        "git_dir_note": (
            "The .git directory is not committed (git would store a gitlink and "
            "clones would get an empty directory). Recreate it with "
            "scripts/make_fixtures.py, or let the pytest `sample_repo` fixture do it."
        ),
        "must_contain": ["class Widget", "widgetlib", "def scaled"],
        "must_not_contain": ["this-should-never-appear-in-ingested-output"],
    }


# ---------------------------------------------------------------------------
# Driver
# ---------------------------------------------------------------------------

#: Ordered because scanned.pdf and corrupt.pdf derive from simple.pdf.
BUILDERS: Final[list[tuple[str, Any]]] = [
    ("simple.pdf", build_simple_pdf),
    ("tables.pdf", build_tables_pdf),
    ("twocolumn.pdf", build_twocolumn_pdf),
    ("scanned.pdf", build_scanned_pdf),
    ("corrupt.pdf", build_corrupt_pdf),
    ("report.docx", build_report_docx),
    ("unicode.docx", build_unicode_docx),
    ("deck.pptx", build_deck_pptx),
    ("data.xlsx", build_data_xlsx),
    ("article.html", build_article_html),
    ("boilerplate.html", build_boilerplate_html),
    ("long_context.md", build_long_context_md),
    ("sample_repo/", build_sample_repo),
]


def generate(out: Path) -> dict[str, Any]:
    """Generate the whole corpus into ``out``.

    Args:
        out: Directory to write fixtures into. Created if absent.

    Returns:
        The ground-truth mapping, keyed by fixture name.
    """
    out.mkdir(parents=True, exist_ok=True)
    ground_truth: dict[str, Any] = {}
    for name, builder in BUILDERS:
        print(f"  building {name}")
        ground_truth[name] = builder(out)

    manifest = {
        "_generated_by": "scripts/make_fixtures.py",
        "_note": (
            "Structural facts about each fixture. Conversion tests assert against "
            "these rather than against byte-identical output."
        ),
        "fixtures": ground_truth,
    }
    (out / "ground_truth.json").write_text(
        json.dumps(manifest, indent=2, ensure_ascii=False, sort_keys=True) + "\n",
        encoding="utf-8",
        newline="\n",
    )
    return ground_truth


def digests(root: Path) -> dict[str, str]:
    """Return SHA-256 digests for every file under ``root``.

    Args:
        root: Directory to walk.

    Returns:
        Mapping of POSIX-style relative path to hex digest, excluding ``.git``.
    """
    result: dict[str, str] = {}
    for path in sorted(root.rglob("*")):
        if not path.is_file() or ".git" in path.relative_to(root).parts:
            continue
        result[path.relative_to(root).as_posix()] = sha256(path)
    return result


def main(argv: list[str] | None = None) -> int:
    """Entry point.

    Args:
        argv: Argument vector; defaults to ``sys.argv[1:]``.

    Returns:
        Process exit code.
    """
    parser = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    parser.add_argument(
        "--out",
        type=Path,
        default=DEFAULT_OUT,
        help=f"output directory (default: {DEFAULT_OUT})",
    )
    parser.add_argument(
        "--check",
        action="store_true",
        help="regenerate into a temp dir and verify the committed corpus is reproducible",
    )
    args = parser.parse_args(argv)

    if args.check:
        print(f"Checking reproducibility against {args.out}")
        with tempfile.TemporaryDirectory() as tmp:
            scratch = Path(tmp) / "fixtures"
            generate(scratch)
            expected = digests(args.out)
            actual = digests(scratch)
        if expected == actual:
            print(f"OK: {len(actual)} files reproduced byte-for-byte")
            return 0
        for name in sorted(set(expected) | set(actual)):
            if expected.get(name) != actual.get(name):
                print(
                    f"MISMATCH {name}: "
                    f"committed={expected.get(name)} regenerated={actual.get(name)}"
                )
        return 1

    print(f"Generating fixtures into {args.out}")
    generate(args.out)
    print(f"Done: {len(digests(args.out))} files")
    return 0


if __name__ == "__main__":
    sys.exit(main())
